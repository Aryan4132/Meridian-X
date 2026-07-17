import os
import re
from typing import Dict, Any, List, Optional
from src.core.audit_logger import log_sensitive_action

# PDF reading & merging
try:
    import pypdf
except ImportError:
    pypdf = None

# PDF generation
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
except ImportError:
    SimpleDocTemplate = None

# Word documents
try:
    import docx
except ImportError:
    docx = None

# Excel spreadsheets
try:
    import openpyxl
except ImportError:
    openpyxl = None

# PowerPoint presentations
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


def read_document_text(file_path: str) -> str:
    """
    Extracts text content or data from office document formats (.pdf, .docx, .pptx, .xlsx, .xls).
    Returns formatted string content.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        if not pypdf:
            raise ImportError("The 'pypdf' package is not installed.")
        reader = pypdf.PdfReader(file_path)
        text_parts = []
        for i, page in enumerate(reader.pages):
            text_parts.append(f"--- Page {i + 1} ---")
            text_parts.append(page.extract_text() or "")
        return "\n".join(text_parts)
        
    elif ext == ".docx":
        if not docx:
            raise ImportError("The 'python-docx' package is not installed.")
        doc = docx.Document(file_path)
        text_parts = []
        for p in doc.paragraphs:
            text_parts.append(p.text)
        for table in doc.tables:
            text_parts.append("\n--- Table ---")
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts)
        
    elif ext == ".pptx":
        if not Presentation:
            raise ImportError("The 'python-pptx' package is not installed.")
        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
        
    elif ext in [".xlsx", ".xls"]:
        if ext == ".xlsx":
            if not openpyxl:
                raise ImportError("The 'openpyxl' package is not installed.")
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                
                # Convert sheet to a Markdown table
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    text_parts.append("(Empty Sheet)")
                    continue
                
                # Filter trailing empty rows/cols to keep the representation clean
                clean_rows = []
                for row in rows:
                    if any(cell is not None for cell in row):
                        clean_rows.append(row)
                
                if not clean_rows:
                    text_parts.append("(Empty Sheet)")
                    continue
                
                # Determine max columns
                max_cols = max(len(r) for r in clean_rows)
                for r in clean_rows:
                    cells = [str(c) if c is not None else "" for c in r]
                    # Pad cells to max_cols
                    cells += [""] * (max_cols - len(cells))
                    text_parts.append(" | " + " | ".join(cells) + " |")
                    
            return "\n".join(text_parts)
        else:
            try:
                import xlrd
            except ImportError:
                raise ImportError("The 'xlrd' package is required to read old Excel .xls files.")
            wb = xlrd.open_workbook(file_path)
            text_parts = []
            for sheet_index in range(wb.nsheets):
                sheet = wb.sheet_by_index(sheet_index)
                text_parts.append(f"--- Sheet: {sheet.name} ---")
                for r in range(sheet.nrows):
                    row_vals = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                    text_parts.append(" | " + " | ".join(row_vals) + " |")
            return "\n".join(text_parts)
            
    else:
        raise ValueError(f"Unsupported document format: {ext}")


def create_word_document(file_path: str, content_markdown: str) -> str:
    """
    Creates a Word (.docx) document from markdown-formatted text.
    Supports headings (# Heading), bullets (- Bullet), and bold/italic inline markdown.
    """
    if not docx:
        raise ImportError("The 'python-docx' package is not installed.")
        
    parent = os.path.dirname(os.path.abspath(file_path))
    if parent:
        os.makedirs(parent, exist_ok=True)
        
    doc = docx.Document()
    
    def parse_inline_and_add(paragraph, text):
        # Parses **bold** and *italic*
        parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                r = paragraph.add_run(part[2:-2])
                r.bold = True
            elif part.startswith("*") and part.endswith("*"):
                r = paragraph.add_run(part[1:-1])
                r.italic = True
            else:
                paragraph.add_run(part)

    for line in content_markdown.splitlines():
        stripped = line.strip()
        if not line: # Preserve empty paragraphs
            doc.add_paragraph()
            continue
            
        if stripped.startswith("# "):
            p = doc.add_paragraph()
            p.style = doc.styles['Heading 1']
            parse_inline_and_add(p, stripped[2:])
        elif stripped.startswith("## "):
            p = doc.add_paragraph()
            p.style = doc.styles['Heading 2']
            parse_inline_and_add(p, stripped[3:])
        elif stripped.startswith("### "):
            p = doc.add_paragraph()
            p.style = doc.styles['Heading 3']
            parse_inline_and_add(p, stripped[4:])
        elif stripped.startswith("- ") or stripped.startswith("* "):
            p = doc.add_paragraph(style='List Bullet')
            parse_inline_and_add(p, stripped[2:])
        else:
            p = doc.add_paragraph()
            parse_inline_and_add(p, line)
            
    doc.save(file_path)
    
    log_sensitive_action(
        category="FILE_WRITE",
        action="create_word_document",
        details={"path": file_path, "markdown_length": len(content_markdown)},
        status="SUCCESS"
    )
    return f"Successfully created Word document at {file_path}"


def edit_word_document(
    file_path: str, 
    action: str, 
    search_text: Optional[str] = None, 
    replace_text: Optional[str] = None, 
    append_markdown: Optional[str] = None
) -> str:
    """
    Edits a Word (.docx) file.
    Actions supported:
    - 'replace_text': Searches and replaces all instances of search_text with replace_text.
    - 'append_text': Appends content formatted as markdown to the end of the document.
    """
    if not docx:
        raise ImportError("The 'python-docx' package is not installed.")
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Word document not found: {file_path}")
        
    doc = docx.Document(file_path)
    
    if action == "replace_text":
        if not search_text:
            raise ValueError("search_text is required for replace_text action")
        replace_text = replace_text or ""
        
        count = 0
        for paragraph in doc.paragraphs:
            if search_text in paragraph.text:
                # To maintain formatting runs as best as possible, we do run replacement or paragraph-level replace
                # Replacing at run level can be fragmented, so we replace paragraph text if it's simpler
                # However, full text replacement at paragraph level resets formatting. We try run-level first:
                # If search_text matches a single run, replace it. Otherwise, replace paragraph text.
                replaced = False
                for run in paragraph.runs:
                    if search_text in run.text:
                        run.text = run.text.replace(search_text, replace_text)
                        replaced = True
                        count += 1
                if not replaced:
                    paragraph.text = paragraph.text.replace(search_text, replace_text)
                    count += 1
                    
        doc.save(file_path)
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_word_document:replace_text",
            details={"path": file_path, "search_text": search_text, "replacements": count},
            status="SUCCESS"
        )
        return f"Successfully replaced {count} instances of '{search_text}' in {file_path}"
        
    elif action == "append_text":
        if not append_markdown:
            raise ValueError("append_markdown is required for append_text action")
            
        def parse_inline_and_add(paragraph, text):
            parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    r = paragraph.add_run(part[2:-2])
                    r.bold = True
                elif part.startswith("*") and part.endswith("*"):
                    r = paragraph.add_run(part[1:-1])
                    r.italic = True
                else:
                    paragraph.add_run(part)

        for line in append_markdown.splitlines():
            stripped = line.strip()
            if not line:
                doc.add_paragraph()
                continue
                
            if stripped.startswith("# "):
                p = doc.add_paragraph()
                p.style = doc.styles['Heading 1']
                parse_inline_and_add(p, stripped[2:])
            elif stripped.startswith("## "):
                p = doc.add_paragraph()
                p.style = doc.styles['Heading 2']
                parse_inline_and_add(p, stripped[3:])
            elif stripped.startswith("### "):
                p = doc.add_paragraph()
                p.style = doc.styles['Heading 3']
                parse_inline_and_add(p, stripped[4:])
            elif stripped.startswith("- ") or stripped.startswith("* "):
                p = doc.add_paragraph(style='List Bullet')
                parse_inline_and_add(p, stripped[2:])
            else:
                p = doc.add_paragraph()
                parse_inline_and_add(p, line)
                
        doc.save(file_path)
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_word_document:append_text",
            details={"path": file_path, "append_length": len(append_markdown)},
            status="SUCCESS"
        )
        return f"Successfully appended markdown content to Word document {file_path}"
        
    else:
        raise ValueError(f"Unsupported action: {action}")


def create_excel_document(file_path: str, sheets_data: Dict[str, List[List[Any]]]) -> str:
    """
    Creates an Excel (.xlsx) document.
    sheets_data is a dictionary where keys are sheet names and values are 2D arrays of cells.
    """
    if not openpyxl:
        raise ImportError("The 'openpyxl' package is not installed.")
        
    parent = os.path.dirname(os.path.abspath(file_path))
    if parent:
        os.makedirs(parent, exist_ok=True)
        
    wb = openpyxl.Workbook()
    # Remove default sheet
    default_sheet = wb.active
    wb.remove(default_sheet)
    
    for sheet_name, rows in sheets_data.items():
        ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
            
    wb.save(file_path)
    log_sensitive_action(
        category="FILE_WRITE",
        action="create_excel_document",
        details={"path": file_path, "sheets": list(sheets_data.keys())},
        status="SUCCESS"
    )
    return f"Successfully created Excel document at {file_path}"


def edit_excel_document(
    file_path: str, 
    action: str, 
    sheet_name: str, 
    range_or_cell: Optional[str] = None, 
    data: Optional[List[List[Any]]] = None, 
    find_text: Optional[str] = None, 
    replace_text: Optional[str] = None
) -> str:
    """
    Edits an Excel (.xlsx) file.
    Actions supported:
    - 'update_cells': Updates a specific cell or range (e.g. 'A1' or 'A1:B2') in sheet_name using data (2D array).
    - 'append_rows': Appends rows (2D array in data) to the end of sheet_name.
    - 'replace_text': Replaces all cells matching find_text with replace_text in sheet_name.
    """
    if not openpyxl:
        raise ImportError("The 'openpyxl' package is not installed.")
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Excel file not found: {file_path}")
        
    wb = openpyxl.load_workbook(file_path)
    if sheet_name not in wb.sheetnames:
        raise ValueError(f"Sheet '{sheet_name}' not found in Excel workbook.")
        
    ws = wb[sheet_name]
    
    if action == "update_cells":
        if not range_or_cell:
            raise ValueError("range_or_cell is required for update_cells action")
        if not data:
            raise ValueError("data is required for update_cells action")
            
        # Check if single cell or range
        if ":" in range_or_cell:
            cells_range = ws[range_or_cell]
            # data should match the dimensions of the range
            # Flatten or match row-by-row
            for i, row in enumerate(cells_range):
                if i >= len(data):
                    break
                for j, cell in enumerate(row):
                    if j >= len(data[i]):
                        break
                    cell.value = data[i][j]
        else:
            ws[range_or_cell] = data[0][0] if isinstance(data[0], list) else data[0]
            
        wb.save(file_path)
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_excel_document:update_cells",
            details={"path": file_path, "sheet": sheet_name, "range": range_or_cell},
            status="SUCCESS"
        )
        return f"Successfully updated cells '{range_or_cell}' in sheet '{sheet_name}' of {file_path}"
        
    elif action == "append_rows":
        if not data:
            raise ValueError("data (list of rows) is required for append_rows action")
            
        for row in data:
            ws.append(row)
            
        wb.save(file_path)
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_excel_document:append_rows",
            details={"path": file_path, "sheet": sheet_name, "rows_count": len(data)},
            status="SUCCESS"
        )
        return f"Successfully appended {len(data)} rows to sheet '{sheet_name}' in {file_path}"
        
    elif action == "replace_text":
        if not find_text:
            raise ValueError("find_text is required for replace_text action")
        replace_text = replace_text or ""
        
        count = 0
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None and str(cell.value) == find_text:
                    cell.value = replace_text
                    count += 1
                    
        wb.save(file_path)
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_excel_document:replace_text",
            details={"path": file_path, "sheet": sheet_name, "replacements": count},
            status="SUCCESS"
        )
        return f"Successfully replaced {count} cells in sheet '{sheet_name}' in {file_path}"
        
    else:
        raise ValueError(f"Unsupported action: {action}")


def create_powerpoint_presentation(file_path: str, slides: List[Dict[str, Any]]) -> str:
    """
    Creates a PowerPoint (.pptx) presentation.
    slides is a list of dicts: [{'title': 'Slide 1 Title', 'content': ['Point A', 'Point B']}, ...]
    """
    if not Presentation:
        raise ImportError("The 'python-pptx' package is not installed.")
        
    parent = os.path.dirname(os.path.abspath(file_path))
    if parent:
        os.makedirs(parent, exist_ok=True)
        
    prs = Presentation()
    
    # 0 is usually Title Slide, 1 is Title and Content layout
    title_content_layout = prs.slide_layouts[1]
    
    for slide_data in slides:
        slide = prs.slides.add_slide(title_content_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "")
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        content = slide_data.get("content", [])
        if isinstance(content, str):
            tf.text = content
        elif isinstance(content, list):
            # First item in bullet frame
            if content:
                tf.text = content[0]
                for item in content[1:]:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                    
    prs.save(file_path)
    log_sensitive_action(
        category="FILE_WRITE",
        action="create_powerpoint_presentation",
        details={"path": file_path, "slides_count": len(slides)},
        status="SUCCESS"
    )
    return f"Successfully created PowerPoint presentation at {file_path}"


def edit_powerpoint_presentation(
    file_path: str, 
    action: str, 
    slide_index: Optional[int] = None, 
    title: Optional[str] = None, 
    content: Optional[List[str]] = None,
    find_text: Optional[str] = None,
    replace_text: Optional[str] = None
) -> str:
    """
    Edits a PowerPoint (.pptx) file.
    Actions supported:
    - 'add_slide': Appends a slide with title and content.
    - 'replace_text': Searches and replaces text inside all shapes.
    """
    if not Presentation:
        raise ImportError("The 'python-pptx' package is not installed.")
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
    prs = Presentation(file_path)
    
    if action == "add_slide":
        title = title or ""
        content = content or []
        
        title_content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_content_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        if content:
            tf.text = content[0]
            for item in content[1:]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                
        prs.save(file_path)
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_powerpoint_presentation:add_slide",
            details={"path": file_path, "title": title},
            status="SUCCESS"
        )
        return f"Successfully added slide titled '{title}' to {file_path}"
        
    elif action == "replace_text":
        if not find_text:
            raise ValueError("find_text is required for replace_text action")
        replace_text = replace_text or ""
        
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and find_text in shape.text:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if find_text in paragraph.text:
                                paragraph.text = paragraph.text.replace(find_text, replace_text)
                                count += 1
                                
        prs.save(file_path)
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_powerpoint_presentation:replace_text",
            details={"path": file_path, "find_text": find_text, "replacements": count},
            status="SUCCESS"
        )
        return f"Successfully replaced {count} occurrences of '{find_text}' in PowerPoint presentation {file_path}"
        
    else:
        raise ValueError(f"Unsupported action: {action}")


def create_pdf_document(file_path: str, content_markdown: str) -> str:
    """
    Generates a PDF document from Markdown text using reportlab.
    Supports headings (# Heading), bullets (- Bullet), and bold/italic markup.
    """
    if SimpleDocTemplate is None:
        raise ImportError("The 'reportlab' package is not installed.")
        
    parent = os.path.dirname(os.path.abspath(file_path))
    if parent:
        os.makedirs(parent, exist_ok=True)
        
    doc = SimpleDocTemplate(file_path, pagesize=letter)
    styles = getSampleStyleSheet()
    
    # Helper to convert basic md to html-like tags for reportlab Paragraph
    def md_to_html(text: str) -> str:
        text = re.sub(r'\*\*(.*?)\*\*|__(.*?)__', r'<b>\1\2</b>', text)
        text = re.sub(r'\*(.*?)\*|_(.*?)_', r'<i>\1\2</i>', text)
        return text

    story = []
    
    for line in content_markdown.splitlines():
        stripped = line.strip()
        if not line:
            story.append(Spacer(1, 10))
            continue
            
        if stripped.startswith("# "):
            title = md_to_html(stripped[2:])
            story.append(Paragraph(title, styles['Title']))
            story.append(Spacer(1, 12))
        elif stripped.startswith("## "):
            h1 = md_to_html(stripped[3:])
            story.append(Paragraph(h1, styles['Heading1']))
            story.append(Spacer(1, 10))
        elif stripped.startswith("### "):
            h2 = md_to_html(stripped[4:])
            story.append(Paragraph(h2, styles['Heading2']))
            story.append(Spacer(1, 8))
        elif stripped.startswith("- ") or stripped.startswith("* "):
            item = md_to_html(stripped[2:])
            story.append(Paragraph(f"&bull; {item}", styles['Normal']))
            story.append(Spacer(1, 4))
        else:
            body = md_to_html(line)
            story.append(Paragraph(body, styles['Normal']))
            story.append(Spacer(1, 6))
            
    doc.build(story)
    
    log_sensitive_action(
        category="FILE_WRITE",
        action="create_pdf_document",
        details={"path": file_path, "markdown_length": len(content_markdown)},
        status="SUCCESS"
    )
    return f"Successfully created PDF document at {file_path}"


def edit_pdf_document(
    file_path: str, 
    action: str, 
    merge_with_path: Optional[str] = None, 
    append_markdown: Optional[str] = None
) -> str:
    """
    Edits a PDF document. Since PDFs are not easily modified in-place, the actions supported are:
    - 'merge': Merges the current PDF file with another PDF file (specified in merge_with_path), saving the output at file_path.
    - 'append_pages': Generates a temporary PDF from append_markdown and appends it to the end of the PDF at file_path.
    """
    if not pypdf:
        raise ImportError("The 'pypdf' package is not installed.")
        
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PDF document not found: {file_path}")
        
    if action == "merge":
        if not merge_with_path:
            raise ValueError("merge_with_path is required for merge action")
        if not os.path.exists(merge_with_path):
            raise FileNotFoundError(f"PDF file to merge not found: {merge_with_path}")
            
        writer = pypdf.PdfWriter()
        writer.append(file_path)
        writer.append(merge_with_path)
        
        # Temp save then move to keep atomic
        temp_path = file_path + ".tmp"
        writer.write(temp_path)
        writer.close()
        
        if os.path.exists(file_path):
            os.remove(file_path)
        os.rename(temp_path, file_path)
        
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_pdf_document:merge",
            details={"path": file_path, "merged_with": merge_with_path},
            status="SUCCESS"
        )
        return f"Successfully merged {merge_with_path} into {file_path}"
        
    elif action == "append_pages":
        if not append_markdown:
            raise ValueError("append_markdown is required for append_pages action")
            
        temp_pdf = file_path + ".append.tmp.pdf"
        try:
            create_pdf_document(temp_pdf, append_markdown)
            
            writer = pypdf.PdfWriter()
            writer.append(file_path)
            writer.append(temp_pdf)
            
            temp_path = file_path + ".tmp"
            writer.write(temp_path)
            writer.close()
            
            if os.path.exists(file_path):
                os.remove(file_path)
            os.rename(temp_path, file_path)
        finally:
            if os.path.exists(temp_pdf):
                os.remove(temp_pdf)
                
        log_sensitive_action(
            category="FILE_WRITE",
            action="edit_pdf_document:append_pages",
            details={"path": file_path, "append_length": len(append_markdown)},
            status="SUCCESS"
        )
        return f"Successfully appended pages generated from markdown to {file_path}"
        
    else:
        raise ValueError(f"Unsupported action: {action}")
